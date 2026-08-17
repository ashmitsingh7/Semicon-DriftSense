#!/usr/bin/env python3
"""Build solution_presentation.pptx from markdown content using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re
import os

# Read the markdown content
markdown_path = os.path.join(os.path.dirname(__file__), 'presentation_content.md')
with open(markdown_path, 'r') as f:
    content = f.read()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0x1A, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF6, 0xFA)
MEDIUM_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)

def add_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, line_spacing=1.15, font_name='Calibri'):
    """Add textbox with multiple lines (list of (text, bold, font_size, color) tuples)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_info, tuple):
            text, bold, fs, cl = line_info
        else:
            text, bold, fs, cl = line_info, False, font_size, color
        p.text = text
        p.font.size = Pt(fs)
        p.font.bold = bold
        p.font.color.rgb = cl
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(fs * (line_spacing - 1) * 4)
    return txBox

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    """Add a formatted table."""
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = 'Calibri'
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_GRAY
                    paragraph.alignment = PP_ALIGN.CENTER
            # Row coloring
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape

def parse_slides(markdown):
    """Parse markdown into slide structures."""
    slides = []
    current_slide = None
    for line in markdown.split('\n'):
        line = line.rstrip()
        # Slide delimiter
        if line.startswith('## Slide'):
            if current_slide:
                slides.append(current_slide)
            # Extract slide number and title
            match = re.match(r'## Slide (\d+):\s*(.*)', line)
            slide_num = int(match.group(1)) if match else len(slides) + 1
            title = match.group(2) if match else ""
            current_slide = {'num': slide_num, 'title': title, 'content': []}
        elif current_slide is not None:
            current_slide['content'].append(line)
    if current_slide:
        slides.append(current_slide)
    return slides

slides_data = parse_slides(content)

# Build each slide
for slide_data in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, WHITE)

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Slide title in header
    add_textbox(slide, 0.5, 0.1, 12, 0.7, f"Slide {slide_data['num']}: {slide_data['title']}",
                font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Parse content
    content_lines = slide_data['content']
    y_pos = 1.1

    i = 0
    while i < len(content_lines):
        line = content_lines[i]

        if not line.strip():
            y_pos += 0.15
            i += 1
            continue

        # Table detection
        if line.strip().startswith('|') and i + 1 < len(content_lines) and content_lines[i+1].strip().startswith('|'):
            # Collect table rows
            table_lines = []
            while i < len(content_lines) and content_lines[i].strip().startswith('|'):
                table_lines.append(content_lines[i])
                i += 1

            # Parse table
            rows = []
            for tline in table_lines:
                cells = [c.strip() for c in tline.split('|') if c.strip()]
                if cells:
                    rows.append(cells)

            if rows:
                n_rows = len(rows)
                n_cols = len(rows[0]) if rows else 0
                col_widths = [3.5, 9.5] if n_cols == 2 else None
                if n_cols > 2:
                    col_widths = [2.0] + [1.3] * (n_cols - 1)
                add_table(slide, 0.5, y_pos, 12.3, 0.35 * n_rows, n_rows, n_cols, rows, col_widths)
                y_pos += 0.35 * n_rows + 0.3
            continue

        # Bullet points
        if line.strip().startswith('- ') or line.strip().startswith('* '):
            bullet_text = line.strip()[2:].strip()
            add_textbox(slide, 0.8, y_pos, 11.5, 0.4, f"• {bullet_text}", font_size=15, color=DARK_GRAY)
            y_pos += 0.35

        # Bold section headers
        elif line.startswith('**') and line.endswith('**'):
            add_textbox(slide, 0.5, y_pos, 12, 0.4, line.strip('*'), font_size=18, bold=True, color=ACCENT_BLUE)
            y_pos += 0.4

        # Regular text
        elif line.startswith('# '):
            add_textbox(slide, 0.5, y_pos, 12, 0.5, line[2:], font_size=26, bold=True, color=DARK_BLUE)
            y_pos += 0.6
        elif line.startswith('### '):
            add_textbox(slide, 0.5, y_pos, 12, 0.4, line[4:], font_size=20, bold=True, color=ACCENT_BLUE)
            y_pos += 0.45
        else:
            # Wrap long lines
            wrapped = line
            add_textbox(slide, 0.5, y_pos, 12.3, 0.4, wrapped, font_size=14, color=DARK_GRAY)
            y_pos += 0.35

        i += 1

        # Prevent overflow
        if y_pos > 6.8:
            break

# Save
output_path = os.path.join(os.path.dirname(__file__), 'solution_presentation.pptx')
prs.save(output_path)
print(f"Saved PPTX to {output_path}")